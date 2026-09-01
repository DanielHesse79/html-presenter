"""Export a <deck-stage> deck to .pptx as a presentation-safe backup.

Each slide becomes a full-bleed image on a matching 16:9 canvas, and the
speaker notes are written into PowerPoint's notes pane so Presenter View
works if the HTML deck ever fails you on stage.

    python tools/export_pptx.py my-deck.html
    python tools/export_pptx.py my-deck.html -o backup.pptx --scale 2

Animation, audio and interactivity do not survive the trip — this is a
static fallback, not a conversion.

Requires: pip install python-pptx PyMuPDF
"""

from __future__ import annotations

import argparse
import json
import re
import sys
import tempfile
from pathlib import Path

sys.path.insert(0, str(Path(__file__).parent))
from _chrome import ChromeNotFound, print_to_pdf  # noqa: E402

EMU_PER_PX = 9525  # 1 px at 96 DPI


def parse_size(value: str) -> tuple[int, int]:
    try:
        w, h = value.lower().split("x")
        return int(w), int(h)
    except ValueError:
        raise argparse.ArgumentTypeError(
            f"--size expects WIDTHxHEIGHT, e.g. 1920x1080 (got {value!r})"
        )


def read_notes(html: str) -> list[str]:
    """Pull the speaker-notes JSON array out of the deck document."""
    match = re.search(
        r'<script[^>]*id=["\']speaker-notes["\'][^>]*>(.*?)</script>',
        html, re.DOTALL,
    )
    if not match:
        return []
    try:
        parsed = json.loads(match.group(1))
    except json.JSONDecodeError as e:
        print(f"warning: could not parse #speaker-notes ({e}) — exporting without notes",
              file=sys.stderr)
        return []
    return parsed if isinstance(parsed, list) else []


def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__,
                                 formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument("deck", type=Path, help="path to the deck .html")
    ap.add_argument("-o", "--out", type=Path, help="output .pptx (default: '<deck> (backup).pptx')")
    ap.add_argument("--size", type=parse_size, default=(1920, 1080),
                    help="slide size in px (default: 1920x1080)")
    ap.add_argument("--scale", type=float, default=2.0,
                    help="raster oversampling for crisper text (default: 2)")
    ap.add_argument("--chrome", help="path to a Chrome/Chromium/Edge binary")
    args = ap.parse_args()

    try:
        import fitz  # PyMuPDF
        from pptx import Presentation
        from pptx.util import Emu
    except ImportError as e:
        print(f"error: missing dependency ({e.name}). Run:\n"
              f"  pip install python-pptx PyMuPDF", file=sys.stderr)
        return 1

    if not args.deck.is_file():
        print(f"error: {args.deck} not found", file=sys.stderr)
        return 1

    out = args.out or args.deck.with_name(f"{args.deck.stem} (backup).pptx")
    width, height = args.size

    print(f"[1/3] Rendering {args.deck.name} at {width}x{height}...")
    with tempfile.TemporaryDirectory(prefix="deck-pptx-") as tmpdir:
        tmp = Path(tmpdir)
        pdf = tmp / "deck.pdf"
        try:
            print_to_pdf(args.deck, pdf, width=width, height=height, chrome=args.chrome)
        except ChromeNotFound as e:
            print(f"error: {e}", file=sys.stderr)
            return 1

        print(f"[2/3] Rasterising at {args.scale}x...")
        doc = fitz.open(pdf)
        matrix = fitz.Matrix(args.scale, args.scale)
        images = []
        for i, page in enumerate(doc):
            image = tmp / f"slide-{i + 1:03d}.png"
            page.get_pixmap(matrix=matrix, alpha=False).save(image)
            images.append(image)
        doc.close()
        print(f"      {len(images)} slide(s)")

        notes = read_notes(args.deck.read_text(encoding="utf-8"))
        if notes and len(notes) != len(images):
            print(f"      note: {len(notes)} note(s) for {len(images)} slide(s) — "
                  f"the shorter list wins", file=sys.stderr)

        print("[3/3] Building .pptx...")
        prs = Presentation()
        prs.slide_width = Emu(width * EMU_PER_PX)
        prs.slide_height = Emu(height * EMU_PER_PX)
        blank = prs.slide_layouts[6]

        for i, image in enumerate(images):
            slide = prs.slides.add_slide(blank)
            slide.shapes.add_picture(
                str(image), Emu(0), Emu(0),
                width=prs.slide_width, height=prs.slide_height,
            )
            if i < len(notes) and notes[i]:
                slide.notes_slide.notes_text_frame.text = notes[i]

        out.parent.mkdir(parents=True, exist_ok=True)
        prs.save(out)

    attached = sum(1 for n in notes[:len(images)] if n)
    print(f"\nWrote {out}  ({out.stat().st_size / 1024:.0f} KB, "
          f"{len(images)} slides, {attached} notes)")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
